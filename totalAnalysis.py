import os
from classes import *
from plotFuncs import *
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from openpyxl.styles import Alignment, borders
from openpyxl.styles.borders import Border, Side

# 1. Loading the file with data for total period
file_path = os.getcwd() + '/data/total.xlsx'
total_label = 'Total'

myWorkbook = MyWorkbook(file_path)
myWorksheets = myWorkbook.mywb.sheetnames
n_of_months = len(myWorkbook.sheets_list)
start_label = [int(myWorksheets[0][:2]), int(myWorksheets[0][-2:])]


# 2. Preparing data from the parsed sheets for visualization
# -- Total as a whole
# a) Barplot for all categories
index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
values_desc = [myWorkbook.cats_sums_list[i] for i in index_order
               if myWorkbook.cats_sums_list[i] > 0]
labels_desc = [myWorkbook.cats_names[i] for i in index_order
               if myWorkbook.cats_sums_list[i] > 0]

# b) Piechart of spendings for the main categories
index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
_top_indices = index_order[0:5]
_low_indices = index_order[5:len(index_order)]

top_labels = [myWorkbook.cats_names[i] for i in _top_indices]
top_values = [myWorkbook.cats_sums_list[i] for i in _top_indices]
others_values = [myWorkbook.cats_sums_list[i] for i in _low_indices]

top_plus_others_values = top_values + [sum(others_values)]
top_plus_others_labels_0 = top_labels + ['Pozostałe']
top_plus_others_labels = [top_plus_others_labels_0[i] + ' - '
                          + str(round(top_plus_others_values[i], 2)) + ' zł'
                          for i in range(len(top_plus_others_values))]

# c) Piechart of the metacategories
metacats_values = [myWorkbook.sum_basic, myWorkbook.sum_addit,
                   myWorkbook.sum_giftdon]
metacats_labels = ['Podstawowe - ' + str(round(myWorkbook.sum_basic, 2))
                   + 'zł', 'Dodatkowe - ' + str(round(myWorkbook.sum_addit, 2))
                   + 'zł', 'Prezenty i donacje - '
                   + str(round(myWorkbook.sum_giftdon, 2)) + 'zł']

# d) Piechart of incomes
_values_list_inc = list(myWorkbook.incomes_dict.values())
_labels_list_inc = list(myWorkbook.incomes_dict.keys())
incomes_values = [inc for inc in _values_list_inc if inc > 0]
incomes_labels = [inc+' - '+str(_values_list_inc[i])+'zł' for i, inc in
                  enumerate(_labels_list_inc) if _values_list_inc[i] > 0]

# -- Averaged month --
# e) Piechart of spendings for the main categories
top_plus_others_values_avg = [i/n_of_months for i in top_plus_others_values]
top_plus_others_labels_avg = [top_plus_others_labels_0[i] + ' - '
                              + str(round(top_plus_others_values[i]
                                          / n_of_months, 2)) + ' zł'
                              for i in range(len(top_plus_others_values))]

# f) Piechart of the metacategories
metacats_values_avg = [i/n_of_months for i in metacats_values]
metacats_labels_avg = ['Podstawowe - '
                       + str(round(myWorkbook.sum_basic/n_of_months, 2))
                       + 'zł', 'Dodatkowe - '
                       + str(round(myWorkbook.sum_addit/n_of_months, 2))
                       + 'zł', 'Prezenty i donacje - '
                       + str(round(myWorkbook.sum_giftdon/n_of_months, 2))
                       + 'zł']

# g) Piechart of incomes
incomes_values_avg = [i/n_of_months for i in incomes_values]
incomes_labels_avg = [inc + ' - ' + str(round(_values_list_inc[i]
                      / n_of_months, 2)) + 'zł'
                      for i, inc in enumerate(_labels_list_inc)
                      if _values_list_inc[i] > 0]

# -- Total as a sequence of months --
# h) Stackplot of cummulated spendings for the top categories
top_spends_seqs = []
# top categories
for c, cat in enumerate(top_labels):
    top_spends_seqs.append([])
    for month in myWorkbook.sheets_list:
        top_spends_seqs[c].append(month.cats_sums[cat])
# others
top_spends_seqs.append([])
for m, month in enumerate(myWorkbook.sheets_list):
    top_spends_seqs[-1].append(np.sum(month.cats_sums_list)
                               - sum(cat[m] for cat in top_spends_seqs[:-1]
                               if cat[m]))
# cummulated sums of the lists
top_spends_seqs_cumsum = []
for seq in top_spends_seqs:
    top_spends_seqs_cumsum.append(np.cumsum(seq))

# i) Lineplot of cummulated spendings, incomes and savings
line_incomes = []
line_spendings = []
line_balance = []
line_savings = []

for month in myWorkbook.sheets_list:
    line_incomes.append(month.incomes)
    line_spendings.append(sum(month.cats_sums_list))
    line_balance.append(month.balance[0])
    line_savings.append(month.savings)

line_incomes = np.cumsum(line_incomes)
line_spendings = np.cumsum(line_spendings)
line_balance = np.cumsum(line_balance)
line_savings = np.cumsum(line_savings)

# j) Lineplot of spendings and incomes in subsequent months
spendings_list = []
incomes_list = []
for month in myWorkbook.sheets_list:
    spendings_list.append(month.sum_total)
    incomes_list.append(month.incomes)

# k) Lineplot of average spendings for subsequent categories so far
current_means_seqs = []
for c, cat in enumerate(top_spends_seqs):
    current_means_seqs.append([])
    for m, month in enumerate(cat):
        current_means_seqs[-1].append(np.mean(top_spends_seqs[c][:m+1]))

# l) Lineplot of main sources
#  choosing of the relevant sources
incomes_main = []
for inc in myWorkbook.incomes_dict:
    if myWorkbook.incomes_dict[inc] > 0.02*myWorkbook.incomes:
        incomes_main.append(inc)

# preprocessing
incomes_seqs = [[] for i in range(len(incomes_main))]
for p, inc in enumerate(incomes_main):
    for m, month in enumerate(myWorkbook.sheets_list):
        try:
            incomes_seqs[p].append(month.incomes_dict[inc])
        except:
            incomes_seqs[p].append(0)


# 3. Visualization and saving the plots
results_dir = os.getcwd() + '/results/' + total_label + ' - wyniki'
if not os.path.exists(results_dir):
    os.mkdir(results_dir)
    os.mkdir(results_dir + '/plots')

# a) Barplot for all categories
values = values_desc
labels = labels_desc
title = total_label + ' - Kwoty wydane w ciągu całego okresu \nna ' \
                      'kolejne kategorie'
fig_name = results_dir + '/plots/plot1.png'

fig = plotBar(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# b) Piechart of spendings for the main categories
values = top_plus_others_values
labels = top_plus_others_labels
title = total_label \
        + ' - Struktura całkowitych wydatków\n z podziałem na kategorie\n\n' \
        'Suma wydatków: ' + str(round(myWorkbook.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot2.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# c) Piechart of the metacategories
values = metacats_values
labels = metacats_labels
title = total_label + ' - Podział wydatków na:\n' \
        + 'Podstawowe, Dodatkowe i Prezenty/Donacje\n\nSuma wydatków: ' \
        + str(round(myWorkbook.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot3.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# d) Piechart of incomes
values = incomes_values
labels = incomes_labels
title = total_label + ' - Podział przychodów na poszczególne źródła\n\n' \
        'Suma przychodów: ' + str(round(myWorkbook.incomes, 2)) + 'zł\n' \
        'Nadwyżka przychodów: ' + str(round(myWorkbook.balance[0], 2)) \
        + 'zł  (' + str(round(100*myWorkbook.balance[0]
                               / myWorkbook.incomes, 2)) \
        + '%)\n'

fig_name = results_dir + '/plots/plot4.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# e) Piechart of spendings for the main categories
values = top_plus_others_values_avg
labels = top_plus_others_labels_avg
title = total_label \
        + ' - Struktura wydatków w uśrednionym miesiącu\n z podziałem na ' \
        'kategorie\n\n'+'Suma wydatków: ' \
        + str(round(myWorkbook.sum_total/n_of_months, 2)) + 'zł'
fig_name = results_dir + '/plots/plot5.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# f) Piechart of metacategories
values = metacats_values_avg
labels = metacats_labels_avg
title = total_label + ' - Podział wydatków na: ' \
        + 'Podstawowe, Dodatkowe\n i Prezenty/Donacje ' \
        'w uśrednionym miesiącu\n\n' \
        'Suma wydatków: ' + str(round(myWorkbook.sum_total/n_of_months, 2)) \
        + 'zł\n'
fig_name = results_dir + '/plots/plot6.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# g) Piechart of incomes
values = incomes_values_avg
labels = incomes_labels_avg
title = total_label + ' - Podział przychodów na poszczególne źródła\n' \
        'w uśrednionym miesiącu\n\n' + 'Suma przychodów: ' \
        + str(round(myWorkbook.incomes/n_of_months, 2)) + 'zł\n' \
        + 'Nadwyżka przychodów: ' \
        + str(round(myWorkbook.balance[0]/n_of_months, 2)) \
        + 'zł  (' + str(round(100*myWorkbook.balance[0] / n_of_months
                               / (myWorkbook.incomes/n_of_months), 2)) + '%)'

fig_name = results_dir + '/plots/plot7.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# h) Stackplot of cummulated spendings for the top categories
values = top_spends_seqs_cumsum
labels = top_labels + ['Pozostałe']
title = total_label + ' - Skumulowane wartości wydatków na\n poszczególne ' \
                      'kategorie na przestrzeni całego okresu'
fig_name = results_dir + '/plots/plot8.png'

fig = plotStack(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# i) Lineplot of cummulated spendings, incomes and savings
values = [line_incomes, line_spendings, line_balance, line_savings]
labels = ['Przychody', 'Wydatki', 'Nadwyżka\nprzychodów',
          'Oszczędności\ndługoterminowe']
title = total_label + ' - Skumulowane wartości przychodów, wydatków \n' \
        'i oszczędności na przestrzeni całego okresu'
fig_name = results_dir + '/plots/plot9.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# j) Lineplot of spendings and incomes in subsequent months
values = [incomes_list, spendings_list]
labels = ['Przychody', 'Wydatki']
title = total_label + ' - Przychody i wydatki w kolejnych miesiącach'
fig_name = results_dir + '/plots/plot10.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# k) Lineplot of average spendings for subsequent categories so far
values = current_means_seqs
labels = top_labels + ['Pozostałe']
title = total_label \
        + ' - Dotychczasowe średnie miesięczne wydatki na \nposzczególne ' \
        'kategorie'
fig_name = results_dir + '/plots/plot11.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# l) Lineplot of the sources
values = incomes_seqs
labels = incomes_main
title = total_label + ' - Kwoty przychodów z najważniejszych \n źrodeł na ' \
        'przestrzeni całego okresu'
fig_name = results_dir + '/plots/plot12.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

plt.close('all')

# 4. Pptx presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]
slides = list()

slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].shapes.title
title.text = total_label + ' - raport finansowy'

# Total as a whole
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '1. Total jako całość'

for i in range(4):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(1.5)
    top = Inches(0.0)
    height = width = Inches(7.5)

    pic_path = 'results/' + total_label + ' - wyniki/plots/plot' \
               + str(i+1) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

# Averaged month
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '2. Uśredniony miesiąc'

for i in range(3):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(1.5)
    top = Inches(0.0)
    height = width = Inches(7.5)

    pic_path = 'results/' + total_label + ' - wyniki/plots/plot' \
               + str(i+1+4) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

# Total as a sequence of months
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '3. Total jako sekwencja miesięcy'


for i in range(5):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(0.0)
    top = Inches(0.1)
    height = Inches(10.5)
    width = Inches(7.0)

    pic_path = 'results/' + total_label + ' - wyniki/plots/plot' \
               + str(i+1+7) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)


prs.save(results_dir + '/' + total_label + ' - raport finansowy.pptx')
