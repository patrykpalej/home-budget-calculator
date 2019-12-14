from pptx.util import Inches
from pptx import Presentation


def create_pptx_presentation(year_num, results_dir):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    slides = list()

    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].shapes.title
    title.text = "20" + year_num + " - raport finansowy"

    # Year as a whole
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "1. Rok jako całość"

    for i in range(6):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Averaged month
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "2. Uśredniony miesiąc"

    for i in range(4):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1 + 6) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Year as a sequence of months
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "3. Rok jako sekwencja miesięcy"

    for i in range(7):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.0)
        top = Inches(0.1)
        height = Inches(10.5)
        width = Inches(7.0)

        pic_path = results_dir + "/plots/plot" + str(i + 1 + 10) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    prs.save(results_dir + "/20" + year_num + " - raport finansowy.pptx")
