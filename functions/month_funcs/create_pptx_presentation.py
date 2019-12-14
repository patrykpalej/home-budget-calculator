from pptx.util import Inches
from pptx import Presentation


def create_pptx_presentation(month_num, year_num, results_dir):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    slides = list()

    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].shapes.title
    subtitle = slides[-1].placeholders[1]
    month_dict = {"01": "Styczeń", "02": "Luty", "03": "Marzec",
                  "04": "Kwiecień", "05": "Maj", "06": "Czerwiec",
                  "07": "Lipiec", "08": "Sierpień", "09": "Wrzesień",
                  "10": "Październik", "11": "Listopad", "12": "Grudzień"}
    title.text = month_dict[month_num] + " 20" + year_num \
        + " - raport finansowy"
    subtitle.text = month_num + ".20" + year_num

    for i in range(8):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    prs.save(results_dir + "/" + month_num + ".20" + year_num
             + " - raport finansowy.pptx")
