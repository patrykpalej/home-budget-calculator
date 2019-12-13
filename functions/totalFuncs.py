import openpyxl
import numpy as np
from pptx.util import Inches
from pptx import Presentation
import matplotlib.pyplot as plt
from openpyxl.styles import Alignment
from openpyxl.styles import PatternFill
from openpyxl.styles.borders import Border, Side

from functions.plotFuncs import plotPie, plotBar, plotLine, plotStack, \
    plotScatter


def create_all_plots(myWorkbook, myWorksheets, total_label, results_dir,
                     start_label, n_of_months):
    # -- Total as a whole
    # a) Barplot for all categories
    index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
    values_desc = [myWorkbook.cats_sums_list[i] for i in index_order
                   if myWorkbook.cats_sums_list[i] > 0]
    labels_desc = [myWorkbook.cats_names[i] for i in index_order
                   if myWorkbook.cats_sums_list[i] > 0]

    values = values_desc
    labels = labels_desc
    title = total_label + " - Kwoty wydane w ciągu całego okresu \nna " \
                          "kolejne kategorie"
    fig_name = results_dir + "/plots/plot1.png"

    fig = plotBar(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # b) Piechart of spendings for the main categories
    index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
    _top_indices = index_order[0:5]
    _low_indices = index_order[5:len(index_order)]

    top_labels = [myWorkbook.cats_names[i] for i in _top_indices]
    top_values = [myWorkbook.cats_sums_list[i] for i in _top_indices]
    others_values = [myWorkbook.cats_sums_list[i] for i in _low_indices]

    top_plus_others_values = top_values + [sum(others_values)]
    top_plus_others_labels_0 = top_labels + ["Pozostałe"]
    top_plus_others_labels = [top_plus_others_labels_0[i] + " - "
                              + str(
        round(top_plus_others_values[i], 2)) + " zł"
                              for i in range(len(top_plus_others_values))]

    values = top_plus_others_values
    labels = top_plus_others_labels
    title = total_label \
            + " - Struktura całkowitych wydatków\n z podziałem na kategorie\n\n" \
              "Suma wydatków: " + str(round(myWorkbook.sum_total, 2)) + "zł\n"
    fig_name = results_dir + "/plots/plot2.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # c) Piechart of the metacategories
    metacats_values = [myWorkbook.sum_basic, myWorkbook.sum_addit,
                       myWorkbook.sum_giftdon]
    metacats_labels = ["Podstawowe - " + str(round(myWorkbook.sum_basic, 2))
                       + "zł",
                       "Dodatkowe - " + str(round(myWorkbook.sum_addit, 2))
                       + "zł", "Prezenty i donacje - "
                       + str(round(myWorkbook.sum_giftdon, 2)) + "zł"]

    values = metacats_values
    labels = metacats_labels
    title = total_label + " - Podział wydatków na:\n" \
            + "Podstawowe, Dodatkowe i Prezenty/Donacje\n\nSuma wydatków: " \
            + str(round(myWorkbook.sum_total, 2)) + "zł\n"
    fig_name = results_dir + "/plots/plot3.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # d) Piechart of incomes
    _values_list_inc = list(myWorkbook.incomes_dict.values())
    _labels_list_inc = list(myWorkbook.incomes_dict.keys())
    # marks which incomes are "others"
    others_markers_inc = [1 if i < 0.02 * sum(_values_list_inc) else 0
                          for i in _values_list_inc]

    incomes_values = [inc for i, inc in enumerate(_values_list_inc)
                      if inc > 0 and others_markers_inc[i] == 0]
    incomes_labels = [inc + " - " + str(_values_list_inc[i]) + "zł" for i, inc
                      in
                      enumerate(_labels_list_inc)
                      if
                      _values_list_inc[i] > 0 and others_markers_inc[i] == 0]

    if sum(others_markers_inc) > 0:
        sum_others = 0
        for i, inc in enumerate(_values_list_inc):
            if others_markers_inc[i] == 1:
                sum_others += inc

        incomes_values.append(sum_others)
        incomes_labels.append("Inne - " + str(incomes_values[-1]) + "zł")

    values = incomes_values
    labels = incomes_labels
    title = total_label + " - Podział przychodów na \nposzczególne źródła\n\n" \
                          "Suma przychodów: " + str(
        round(myWorkbook.incomes, 2)) + "zł\n" \
                                        "Nadwyżka przychodów: " + str(
        round(myWorkbook.balance[0], 2)) \
            + "zł  (" + str(round(100 * myWorkbook.balance[0]
                                  / myWorkbook.incomes, 2)) \
            + "%)\n"

    fig_name = results_dir + "/plots/plot4.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # e) Piechart of food subcategories
    amounts = myWorkbook.spends_values_yr["Jedzenie"]
    subcats = myWorkbook.spends_items_yr["Jedzenie"]
    subcats_dict = {}

    for i, subcat in enumerate(subcats):
        if subcat in list(subcats_dict.keys()):
            subcats_dict[subcat] += amounts[i]
        else:
            subcats_dict[subcat] = amounts[i]

    subcats_values = list(subcats_dict.values())
    subcats_labels = [list(subcats_dict.keys())[i] + " - "
                      + str(round(list(subcats_dict.values())[i], 2)) + "zł"
                      for i, sc in enumerate(list(subcats_dict.keys()))]
    # ----
    _subcats_fractions = [sc / sum(subcats_values) for sc in subcats_values]
    subcats_values_with_others = []
    subcats_labels_with_others = []
    others_sum = 0
    for i, sc_f in enumerate(_subcats_fractions):
        if sc_f > 0.02:
            subcats_values_with_others.append(subcats_values[i])
            subcats_labels_with_others.append(subcats_labels[i])
        else:
            others_sum += subcats_values[i]

    if others_sum > 0:
        subcats_values_with_others.append(others_sum)
        subcats_labels_with_others.append("inne - " + str(others_sum) + "zł")

    values = subcats_values_with_others
    labels = subcats_labels_with_others
    title = total_label + " - Podział wydatków spożywczych\n\nCałkowita kwota: " \
            + str(round(myWorkbook.cats_sums["Jedzenie"], 2)) + " zł\n"
    fig_name = results_dir + "/plots/plot5.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # -- Averaged month --
    # f) Piechart of spendings for the main categories
    top_plus_others_values_avg = [i / n_of_months for i in
                                  top_plus_others_values]
    top_plus_others_labels_avg = [top_plus_others_labels_0[i] + " - "
                                  + str(round(top_plus_others_values[i]
                                              / n_of_months, 2)) + " zł"
                                  for i in range(len(top_plus_others_values))]

    values = top_plus_others_values_avg
    labels = top_plus_others_labels_avg
    title = total_label \
            + " - Struktura wydatków w uśrednionym \nmiesiącu z podziałem na " \
              "kategorie\n\n" + "Suma wydatków: " \
            + str(round(myWorkbook.sum_total / n_of_months, 2)) + "zł"
    fig_name = results_dir + "/plots/plot6.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # g) Piechart of the metacategories
    metacats_values_avg = [i / n_of_months for i in metacats_values]
    metacats_labels_avg = ["Podstawowe - "
                           + str(round(myWorkbook.sum_basic / n_of_months, 2))
                           + "zł", "Dodatkowe - "
                           + str(round(myWorkbook.sum_addit / n_of_months, 2))
                           + "zł", "Prezenty i donacje - "
                           + str(
        round(myWorkbook.sum_giftdon / n_of_months, 2))
                           + "zł"]

    values = metacats_values_avg
    labels = metacats_labels_avg
    title = total_label + " - Podział wydatków na: " \
            + "Podstawowe, \nDodatkowe i Prezenty/Donacje " \
              "w uśrednionym miesiącu\n\n" \
              "Suma wydatków: " + str(
        round(myWorkbook.sum_total / n_of_months, 2)) \
            + "zł\n"
    fig_name = results_dir + "/plots/plot7.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # h) Piechart of incomes
    incomes_values_avg = [i / len(myWorksheets) for i in incomes_values]
    incomes_labels_avg = [inc + " - " + str(round(_values_list_inc[i]
                                                  / len(myWorksheets),
                                                  2)) + "zł"
                          for i, inc in enumerate(_labels_list_inc)
                          if others_markers_inc[i] == 0]
    if sum(others_markers_inc) > 0:
        incomes_labels_avg.append("Inne - " + str(round(incomes_values[-1]
                                                        / len(myWorksheets),
                                                        2))
                                  + "zł")

    values = incomes_values_avg
    labels = incomes_labels_avg
    title = total_label + " - Podział przychodów na poszczególne \nźródła" \
                          "w uśrednionym miesiącu\n\n" + "Suma przychodów: " \
            + str(round(myWorkbook.incomes / n_of_months, 2)) + "zł\n" \
            + "Nadwyżka przychodów: " \
            + str(round(myWorkbook.balance[0] / n_of_months, 2)) \
            + "zł  (" + str(round(100 * myWorkbook.balance[0] / n_of_months
                                  / (myWorkbook.incomes / n_of_months),
                                  2)) + "%)"

    fig_name = results_dir + "/plots/plot8.png"

    fig = plotPie(values, labels, title)
    plt.savefig(figure=fig, fname=fig_name)

    # -- Total as a sequence of months --
    # i) Stackplot of cummulated spendings for the top categories
    top_spends_seqs = []
    # top categories
    for c, cat in enumerate(top_labels):
        top_spends_seqs.append([])
        for month in myWorkbook.sheets_list:
            top_spends_seqs[c].append(month.cats_sums[cat])
    # others
    top_spends_seqs.append([])
    for m, month in enumerate(myWorkbook.sheets_list):
        top_spends_seqs[-1].append(np.sum(month.cats_sums_list)
                                   - sum(cat[m] for cat in top_spends_seqs[:-1]
                                         if cat[m]))
    # cummulated sums of the lists
    top_spends_seqs_cumsum = []
    for seq in top_spends_seqs:
        top_spends_seqs_cumsum.append(np.cumsum(seq))

    values = top_spends_seqs_cumsum
    labels = top_labels + ["Pozostałe"]
    title = total_label + " - Skumulowane wartości wydatków na\n poszczególne " \
                          "kategorie na przestrzeni całego okresu"
    fig_name = results_dir + "/plots/plot9.png"

    fig = plotStack(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    # j) Lineplot of cummulated spendings, incomes and savings
    line_incomes = []
    line_spendings = []
    line_balance = []
    line_savings = []

    for month in myWorkbook.sheets_list:
        line_incomes.append(month.incomes)
        line_spendings.append(sum(month.cats_sums_list))
        line_balance.append(month.balance[0])
        line_savings.append(month.savings)

    line_incomes = np.cumsum(line_incomes)
    line_spendings = np.cumsum(line_spendings)
    line_balance = np.cumsum(line_balance)
    line_savings = np.cumsum(line_savings)

    values = [line_incomes, line_spendings, line_balance, line_savings]
    labels = ["Przychody", "Wydatki", "Nadwyżka\nprzychodów",
              "Oszczędności\ndługoterminowe"]
    title = total_label + " - Skumulowane wartości przychodów, wydatków \n" \
                          "i oszczędności na przestrzeni całego okresu"
    fig_name = results_dir + "/plots/plot10.png"

    fig = plotLine(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    # k) Lineplot of spendings and incomes in subsequent months
    spendings_list = []
    incomes_list = []
    earnings_list = []  # for the need of later xlsx summary
    surplus_list = []  # for the need of later xlsx summary
    for month in myWorkbook.sheets_list:
        spendings_list.append(month.sum_total)
        incomes_list.append(month.incomes)
        earnings_list.append(month.earnings)
        surplus_list.append(month.incomes - month.sum_total)

    values = [incomes_list, spendings_list]
    labels = ["Przychody", "Wydatki"]
    title = total_label + " - Przychody i wydatki w kolejnych miesiącach"
    fig_name = results_dir + "/plots/plot11.png"

    fig = plotLine(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    # l) Lineplot of average spendings for subsequent categories so far
    current_means_seqs = []
    for c, cat in enumerate(top_spends_seqs):
        current_means_seqs.append([])
        for m, month in enumerate(cat):
            current_means_seqs[-1].append(np.mean(top_spends_seqs[c][:m + 1]))

    values = current_means_seqs
    labels = top_labels + ["Pozostałe"]
    title = total_label \
            + " - Dotychczasowe średnie miesięczne wydatki na \nposzczególne " \
              "kategorie"
    fig_name = results_dir + "/plots/plot12.png"

    fig = plotLine(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    # m) Lineplot of main sources
    #  choosing of the relevant sources
    incomes_main = []
    for inc in myWorkbook.incomes_dict:
        if myWorkbook.incomes_dict[inc] > 0.02 * myWorkbook.incomes:
            incomes_main.append(inc)

    # preprocessing
    incomes_seqs = [[] for i in range(len(incomes_main))]
    for p, inc in enumerate(incomes_main):
        for m, month in enumerate(myWorkbook.sheets_list):
            try:
                incomes_seqs[p].append(month.incomes_dict[inc])
            except:
                incomes_seqs[p].append(0)

    values = incomes_seqs
    labels = incomes_main
    title = total_label + " - Kwoty przychodów z najważniejszych \n źrodeł na " \
                          "przestrzeni całego okresu"
    fig_name = results_dir + "/plots/plot13.png"

    fig = plotLine(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    # n) Scatterplot incomes vs. spendings
    scatter_incomes = [sheet.incomes for sheet in myWorkbook.sheets_list]
    scatter_spendings = [sheet.sum_total for sheet in myWorkbook.sheets_list]

    values = [scatter_incomes, scatter_spendings]
    title = total_label + " - Przychody vs. wydatki"
    fig_name = results_dir + "/plots/plot14.png"

    fig = plotScatter(values, title)
    plt.savefig(figure=fig, fname=fig_name)

    # o) Lineplot of basic and additional spendings
    basic_array = [s.sum_basic for s in myWorkbook.sheets_list]
    addit_array = [s.sum_addit for s in myWorkbook.sheets_list]
    giftdon_array = [s.sum_giftdon for s in myWorkbook.sheets_list]

    values = [basic_array, addit_array, giftdon_array]
    labels = ["Wydatki\npodstawowe", "Wydatki\ndodatkowe",
              "Prezenty\ni donacje"]
    title = total_label + " - Metakategorie wydatków na przestrzeni czasu"
    fig_name = results_dir + "/plots/plot15.png"

    fig = plotLine(values, labels, title, start_label)
    plt.savefig(figure=fig, fname=fig_name)

    plt.close("all")

    return spendings_list, incomes_list, earnings_list, surplus_list


def create_pptx_presentation(total_label, results_dir):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    slides = list()

    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].shapes.title
    title.text = total_label + " - raport finansowy"

    # Total as a whole
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "1. Total jako całość"

    for i in range(5):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = "results/" + total_label + " - wyniki/plots/plot" \
                   + str(i + 1) + ".png"
        pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Averaged month
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "2. Uśredniony miesiąc"

    for i in range(3):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = "results/" + total_label + " - wyniki/plots/plot" \
                   + str(i + 1 + 5) + ".png"
        pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Total as a sequence of months
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "3. Total jako sekwencja miesięcy"

    for i in range(7):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.0)
        top = Inches(0.1)
        height = Inches(10.5)
        width = Inches(7.0)

        pic_path = "results/" + total_label + " - wyniki/plots/plot" \
                   + str(i + 1 + 8) + ".png"
        pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    prs.save(results_dir + "/" + total_label + " - raport finansowy.pptx")


def create_xlsx_report(results_dir, total_label, spendings_list, incomes_list,
                       earnings_list, surplus_list, myWorkbook, start_label,
                       n_of_months):

    summary_wb = openpyxl.Workbook()

    def export_to_excel(cat_name, num_of_ws):
        # 1. Summary of the period
        # a) first table
        ws = summary_wb.active
        ws.title = "Ogólne"

        ws.cell(1, 2).value = "Wydatki [zł]:"
        ws.cell(1, 3).value = "Zarobki [zł]:"
        ws.cell(1, 4).value = "Przychody [zł]:"
        ws.cell(1, 5).value = "Nadwyżki [zł]:"

        ws.cell(2, 1).value = "Średnia: "
        ws.cell(3, 1).value = "Mediana: "
        ws.cell(4, 1).value = "Std: "

        spendings_mms = [np.mean(spendings_list), np.median(spendings_list),
                         np.std(spendings_list)]
        earnings_mms = [np.mean(earnings_list), np.median(earnings_list),
                        np.std(earnings_list)]
        incomes_mms = [np.mean(incomes_list), np.median(incomes_list),
                       np.std(incomes_list)]
        surplus_mms = [np.mean(surplus_list), np.median(surplus_list),
                       np.std(surplus_list)]

        ws.cell(2, 2).value = round(spendings_mms[0], 2)
        ws.cell(3, 2).value = round(spendings_mms[1], 2)
        ws.cell(4, 2).value = round(spendings_mms[2], 2)

        ws.cell(2, 3).value = round(earnings_mms[0], 2)
        ws.cell(3, 3).value = round(earnings_mms[1], 2)
        ws.cell(4, 3).value = round(earnings_mms[2], 2)

        ws.cell(2, 4).value = round(incomes_mms[0], 2)
        ws.cell(3, 4).value = round(incomes_mms[1], 2)
        ws.cell(4, 4).value = round(incomes_mms[2], 2)

        ws.cell(2, 5).value = round(surplus_mms[0], 2)
        ws.cell(3, 5).value = round(surplus_mms[1], 2)
        ws.cell(4, 5).value = round(surplus_mms[2], 2)

        ws.column_dimensions["B"].width = 14
        ws.column_dimensions["C"].width = 14
        ws.column_dimensions["D"].width = 14
        ws.column_dimensions["E"].width = 14
        ws.column_dimensions["H"].width = 10
        ws.column_dimensions["I"].width = 18

        for r in range(2, 6):
            ws.cell(r, 1).alignment = Alignment(horizontal="right")

        for c in range(2, 6):
            ws.cell(1, c).alignment = Alignment(horizontal="center")

        thin_border = Border(left=Side(style="thin"),
                             right=Side(style="thin"),
                             top=Side(style="thin"),
                             bottom=Side(style="thin"))

        for r in range(4):
            for c in range(5):
                ws.cell(r + 1, c + 1).border = thin_border

        # b) second table
        ws.merge_cells(start_row=1, start_column=8, end_row=1, end_column=9)
        ws.cell(1, 8).value = "Całkowita nadwyżka przychodów:"
        ws.cell(2, 8).value = "[zł]"
        ws.cell(2, 9).value = "[% przychodów]"
        ws.cell(3, 8).value = sum(surplus_list)
        ws.cell(3, 9).value = round(
            sum(surplus_list) / sum(incomes_list) * 100, 2)

        for r in range(1, 4):
            for c in range(8, 10):
                ws.cell(r, c).border = thin_border
                ws.cell(r, c).alignment = Alignment(horizontal="center")

        # 2. Listing the spendings
        values = myWorkbook.spends_values_yr[cat_name]
        items = myWorkbook.spends_items_yr[cat_name]
        monthlabels = myWorkbook.spends_monthlabel_yr[cat_name]

        mdict = {1: "Styczeń", 2: "Luty", 3: "Marzec", 4: "Kwiecień", 5: "Maj",
                 6: "Czerwiec", 7: "Lipiec", 8: "Sierpień", 9: "Wrzesień",
                 10: "Październik", 11: "Listopad", 12: "Grudzień"}

        values_dict = dict()
        items_dict = dict()

        for month_num in np.unique(monthlabels):
            _one_month_list = [values[i] for i, j in enumerate(monthlabels)
                               if j == month_num]
            values_dict[month_num] = _one_month_list

            _one_month_list = [items[i] for i, j in enumerate(monthlabels)
                               if j == month_num]
            items_dict[month_num] = _one_month_list

        ws = summary_wb.create_sheet(cat_name, num_of_ws)
        try:
            ws.column_dimensions["A"].width = max([len(i) for i in items]) + 2
        except:
            return summary_wb

        row = 0
        month = start_label[0]
        year = start_label[1]
        for number in range(1, n_of_months + 1):
            if number in monthlabels:
                row += 1
                ws.cell(row, 1).value = \
                    mdict[month] + " " + "20" + str(year) + " - " \
                    + str(sum(values_dict[number])) + "zł"
                ws.merge_cells(start_row=row, start_column=1, end_row=row,
                               end_column=2)
                ws.cell(row, 1).alignment = Alignment(horizontal="center")
                ws.cell(row, 1).fill = PatternFill(fgColor="93e1e6",
                                                   fill_type="solid")

                row += 1
                ws.cell(row, 1).value = "Co?"
                ws.cell(row, 2).value = "Kwota [zł]"
                ws.cell(row, 1).alignment = Alignment(horizontal="center")
                ws.cell(row, 2).alignment = Alignment(horizontal="center")
                ws.cell(row, 1).fill = PatternFill(fgColor="daf6f7",
                                                   fill_type="solid")
                ws.cell(row, 2).fill = PatternFill(fgColor="daf6f7",
                                                   fill_type="solid")

                for i, val in enumerate(values_dict[number]):
                    row += 1
                    ws.cell(row, 1).value = items_dict[number][i]
                    ws.cell(row, 2).value = val

            month += 1
            if month > 12:
                month -= 12
                year += 1

        for r in range(row):
            ws.cell(row=r + 1, column=1).border = thin_border
            ws.cell(row=r + 1, column=2).border = thin_border

        # 2a) Summary of the category (mean, median, std)
        # For monthly sums
        cat_spends_list = [sheet.cats_sums[cat_name]
                           for sheet in myWorkbook.sheets_list]

        ws.merge_cells(start_row=1, start_column=7, end_row=1,
                       end_column=9)
        ws.cell(1, 7).value = "Dla sum miesięcznych:"
        ws.cell(2, 7).value = "Średnia [zł]:"
        ws.cell(3, 7).value = round(np.mean(cat_spends_list), 2)
        ws.column_dimensions["G"].width = 14
        ws.cell(2, 8).value = "Mediana [zł]:"
        ws.cell(3, 8).value = round(np.median(cat_spends_list), 2)
        ws.column_dimensions["H"].width = 14
        ws.cell(2, 9).value = "Std [zł]:"
        ws.cell(3, 9).value = round(np.std(cat_spends_list), 2)

        for r in range(1, 4):
            for c in range(7, 10):
                ws.cell(r, c).alignment = Alignment(horizontal="center")
                ws.cell(r, c).border = thin_border

        # For individual spendings
        all_values = [v for sublist in values_dict.values() for v in sublist]

        ws.merge_cells(start_row=7, start_column=7, end_row=7,
                       end_column=9)
        ws.cell(7, 7).value = "Dla wydatków indywidualnych:"
        ws.cell(8, 7).value = "Średnia [zł]:"
        ws.cell(9, 7).value = round(np.mean(all_values), 2)
        ws.cell(8, 8).value = "Mediana [zł]:"
        ws.cell(9, 8).value = round(np.median(all_values), 2)
        ws.cell(8, 9).value = "Std [zł]:"
        ws.cell(9, 9).value = round(np.std(all_values), 2)

        for r in range(7, 10):
            for c in range(7, 10):
                ws.cell(r, c).alignment = Alignment(horizontal="center")
                ws.cell(r, c).border = thin_border

        return summary_wb

    summary_wb = export_to_excel("Rzeczy i sprzęty", 1)
    summary_wb = export_to_excel("Hobby i przyjemności", 2)
    summary_wb = export_to_excel("Transport i noclegi", 3)
    summary_wb = export_to_excel("Podróże", 4)
    summary_wb = export_to_excel("Abonamenty i usługi", 5)
    summary_wb = export_to_excel("Leki i zdrowie", 6)
    summary_wb = export_to_excel("Książki i nauka", 7)

    summary_wb.save(results_dir + "/" + total_label + " - podsumowanie.xlsx")

