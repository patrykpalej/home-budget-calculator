from pptx.util import Inches
from pptx import Presentation


def create_pptx_presentation(part_label, results_dir, plot_numbers_list):

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    slides = list()

    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].shapes.title
    title.text = part_label + " - raport finansowy"

    # Part as a whole
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "1. Cały przedział"

    for i in range(plot_numbers_list[0]):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.5)
        top = Inches(0.0)
        height = Inches(7.5)
        width = Inches(9.0)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, width, height)

    # Averaged month
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "2. Uśredniony miesiąc"

    for i in range(plot_numbers_list[0], plot_numbers_list[1]):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.5)
        top = Inches(0.0)
        height = Inches(7.5)
        width = Inches(9.0)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, width, height)

    # Part as a sequence of months
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "3. Sekwencja miesięcy"

    for i in range(plot_numbers_list[1], plot_numbers_list[2]):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.0)
        top = Inches(0.1)
        height = Inches(7.0)
        width = Inches(10.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, width, height)

    prs.save(results_dir + "/" + part_label + " - raport finansowy.pptx")
