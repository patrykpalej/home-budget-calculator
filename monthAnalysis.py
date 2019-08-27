from classes import *
from plotFuncs import *
import os
import sys
import numpy as np
from pptx import Presentation
from pptx.util import Inches

# 1. Loading the file with data for one month
#file_path = os.getcwd() + '/data/monthly/2099.02.xlsx' # hardcoded if no args
file_path = os.getcwd()+'/data/monthly/20'+sys.argv[2]+'.'+sys.argv[1]+'.xlsx'
#month_label = '01.2099' # hardcoded if no args
month_label = sys.argv[1] + '.20' + sys.argv[2]
myWorkbook = MyWorkbook(file_path)
myWorksheet = myWorkbook.sheets_list[0]


# 2. Preparing data from the parsed sheet for visualization
# a) Barplot for all categories
index_order = np.flip(np.argsort(myWorksheet.cats_sums_list), axis=0)
values_desc = [myWorksheet.cats_sums_list[i] for i in index_order
               if myWorksheet.cats_sums_list[i] > 0]
labels_desc = [myWorksheet.cats_names[i] for i in index_order
               if myWorksheet.cats_sums_list[i] > 0]

# b) Piechart of spendings for the main categories
index_order = np.flip(np.argsort(myWorksheet.cats_sums_list), axis=0)
_top_indices = index_order[0:5]
_low_indices = index_order[5:len(index_order)]

top_labels = [myWorksheet.cats_names[i] for i in _top_indices]
top_values = [myWorksheet.cats_sums_list[i] for i in _top_indices]
others_values = [myWorksheet.cats_sums_list[i] for i in _low_indices]

top_plus_others_values = top_values + [sum(others_values)]
top_plus_others_labels = top_labels + ['Pozostałe']
top_plus_others_labels = [top_plus_others_labels[i] + ' - '
                          + str(round(top_plus_others_values[i], 2)) + ' zł'
                          for i in range(len(top_plus_others_values))]

# c) Piechart of the metacategories
metacats_values = [myWorksheet.sum_basic, myWorksheet.sum_addit,
                   myWorksheet.sum_giftdon]
metacats_labels = ['Podstawowe - ' + str(round(myWorksheet.sum_basic, 2)) + 'zł',
                   'Dodatkowe - ' + str(round(myWorksheet.sum_addit, 2)) + 'zł',
                   'Prezenty i donacje - '
                   + str(round(myWorksheet.sum_giftdon, 2)) + 'zł']

# d) Piechart of incomes
_values_list = list(myWorksheet.incomes_dict.values())
_labels_list = list(myWorksheet.incomes_dict.keys())
incomes_values = [inc for inc in _values_list if inc > 0]
incomes_labels = [inc + ' - ' + str(_values_list[i])+'zł' for i, inc in
                  enumerate(_labels_list) if _values_list[i] > 0]

# e) Piechart of earnings
_values_list = list(myWorksheet.earnings_dict.values())
_labels_list = list(myWorksheet.earnings_dict.keys())
earnings_values = [ear for ear in _values_list if ear > 0]
earnings_labels = [ear + ' - ' + str(_values_list[i])+'zł' for i, ear in
                   enumerate(_labels_list) if _values_list[i] > 0]

# f) Piechart of food subcategories
amounts = myWorksheet.spends_values['Jedzenie']
subcats = myWorksheet.spends_items['Jedzenie']
subcats_dict = {}

for i, subcat in enumerate(subcats):
    if subcat in list(subcats_dict.keys()):
        subcats_dict[subcat] += amounts[i]
    else:
        subcats_dict[subcat] = amounts[i]

subcats_values = list(subcats_dict.values())
subcats_labels = [list(subcats_dict.keys())[i] + ' - '
                  + str(round(list(subcats_dict.values())[i], 2)) + 'zł'
                  for i, sc in enumerate(list(subcats_dict.keys()))]

_subcats_fractions = [sc/sum(subcats_values) for sc in subcats_values]
subcats_values_with_others = []
subcats_labels_with_others = []
others_sum = 0
for i, sc_f in enumerate(_subcats_fractions):
    if sc_f > 0.02:
        subcats_values_with_others.append(subcats_values[i])
        subcats_labels_with_others.append(subcats_labels[i])
    else:
        others_sum += subcats_values[i]

if others_sum > 0:
    subcats_values_with_others.append(others_sum)
    subcats_labels_with_others.append('inne - '+str(others_sum)+'zł')

# g) Piechart of 'Hobby i przyjemności' items
amounts = myWorksheet.spends_values['Hobby i przyjemności']
subcats = myWorksheet.spends_items['Hobby i przyjemności']
subcats_dict = {}

for i, subcat in enumerate(subcats):
    if subcat in list(subcats_dict.keys()):
        subcats_dict[subcat] += amounts[i]
    else:
        subcats_dict[subcat] = amounts[i]

subcats_values_1 = list(subcats_dict.values())
subcats_labels_1 = [list(subcats_dict.keys())[i] + ' - '
                    + str(round(list(subcats_dict.values())[i], 2)) + 'zł'
                    for i in range(len(list(subcats_dict.keys())))]

# h) Piechart of 'Rzeczy i sprzęty' items
amounts = myWorksheet.spends_values['Rzeczy i sprzęty']
subcats = myWorksheet.spends_items['Rzeczy i sprzęty']
subcats_dict = {}

for i, subcat in enumerate(subcats):
    if subcat in list(subcats_dict.keys()):
        subcats_dict[subcat] += amounts[i]
    else:
        subcats_dict[subcat] = amounts[i]

subcats_values_2 = list(subcats_dict.values())
subcats_labels_2 = [list(subcats_dict.keys())[i]+' - '
                    + str(round(list(subcats_dict.values())[i], 2)) + 'zł'
                    for i in range(len(list(subcats_dict.keys())))]

# 3. Visualization and saving the plots
results_dir = os.getcwd() + '/results/' + month_label + ' - wyniki'
if not os.path.exists(results_dir):
    os.mkdir(results_dir)
    os.mkdir(results_dir + '/plots/')

# a) Barplot for all categories
values = values_desc
labels = labels_desc
title = month_label + ' - Kwoty wydawane miesięcznie \n na kolejne kategorie\n'
fig_name = results_dir + '/plots/plot1.png'

fig = plotBar(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# b) Piechart of spendings for the main categories
values = top_plus_others_values
labels = top_plus_others_labels
title = month_label + \
        ' - Struktura miesięcznych wydatków\n z podziałem na kategorie\n\n' \
        + 'Suma wydatków: ' + str(round(myWorksheet.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot2.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# c) Piechart of the metacategories
values = metacats_values
labels = metacats_labels
title = month_label + ' - Podział wydatków na: \n' \
        + 'Podstawowe, Dodatkowe i Prezenty/Donacje\n\n' + 'Suma wydatków: ' \
        + str(round(myWorksheet.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot3.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# d) Piechart of incomes
values = incomes_values
labels = incomes_labels
title = month_label + ' - Podział przychodów na poszczególne źródła\n\n' \
        + 'Suma przychodów: ' + str(myWorksheet.incomes) + 'zł\n' \
        + 'Nadwyżka przychodów: ' + str(round(myWorksheet.ballance[0], 2)) \
        + 'zł (' + str(round(100*myWorksheet.ballance[0] /
        myWorksheet.incomes, 2)) + '%)\n'
fig_name = results_dir + '/plots/plot4.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# e) Piechart of earnings
values = earnings_values
labels = earnings_labels
title = month_label + ' - Podział zarobków na poszczególne źródła\n\n' \
        + 'Suma zarobków: ' + str(myWorksheet.earnings) + 'zł\n' \
        + 'Nadwyżka zarobków: ' + str(round(myWorksheet.ballance[1], 2)) \
        + 'zł (' + str(round(100*myWorksheet.ballance[1] /
        myWorksheet.earnings, 2)) + '%)\n'
fig_name = results_dir + '/plots/plot5.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# f) Piechart of food subcategories
values = subcats_values_with_others
labels = subcats_labels_with_others
title = month_label + ' - Podział wydatków spożywczych\n\nCałkowita kwota: ' \
        + str(myWorksheet.cats_sums['Jedzenie']) + ' zł\n'
fig_name = results_dir + '/plots/plot6.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# g) Piechart of 'Hobby i przyjemności' items
values = subcats_values_1
labels = subcats_labels_1
title = month_label + ' - Podział wydatków kategorii\n Hobby i przyjemności\n\n'\
        + 'Całkowita kwota: ' \
        + str(myWorksheet.cats_sums['Hobby i przyjemności']) + ' zł\n'
fig_name = results_dir + '/plots/plot7.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# h) Piechart of 'Rzeczy i sprzęty' items
values = subcats_values_2
labels = subcats_labels_2
title = month_label + ' - Podział wydatków kategorii\n Rzeczy i sprzęty\n\n'\
        'Całkowita kwota: ' \
        + str(myWorksheet.cats_sums['Rzeczy i sprzęty']) + ' zł\n'
fig_name = results_dir + '/plots/plot8.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

plt.close('all')

# 4. Pptx presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]
slides = list()

slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].shapes.title
subtitle = slides[-1].placeholders[1]
month_dict = {'01': 'Styczeń', '02': 'Luty', '03': 'Marzec',
              '04': 'Kwiecień', '05': 'Maj', '06': 'Czerwiec',
              '07': 'Lipiec', '08': 'Sierpień', '09': 'Wrzesień',
              '10': 'Październik', '11': 'Listopad', '12': 'Grudzień'}
title.text = month_dict[sys.argv[1]]+' 20' + sys.argv[2] + ' - raport finansowy'
subtitle.text = sys.argv[1]+'.20'+sys.argv[2]

for i in range(8):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(1.5)
    top = Inches(0.0)
    height = width = Inches(7.5)

    pic_path = results_dir + '/plots/plot' + str(i+1) + '.png'
    slides[-1].shapes.add_picture(pic_path, left, top, height, width)

prs.save(results_dir+'/20' + sys.argv[2] + '.' + sys.argv[1]
         + ' - raport finansowy.pptx')
